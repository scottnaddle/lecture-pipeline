#!/usr/bin/env python3
"""
Phase 1: 슬라이드 텍스트 + 원고 텍스트 추출.
- 슬라이드(.pptx) → slide_texts.json (헤더/푸터/제목/본문 분리)
- 원고(.md) → manuscript_texts.json (H1/H2/단락 분리)
- 두 데이터를 챕터 단위로 묶어 chapter_texts.json 생성
"""
from __future__ import annotations
import json
import re
import unicodedata
from pathlib import Path
from pptx import Presentation

BASE = Path("/Users/scott/Library/CloudStorage/GoogleDrive-scott@naddle.net/내 드라이브/개인/개인 자료/[2024] 대구교대 강의/26년2학기")
MANIFEST = BASE / "강의영상_제작" / "00_매니페스트" / "manifest.json"
OUT = BASE / "강의영상_제작" / "01_스크립트"

# 헤더/푸터로 인식되어 본문 추출에서 제외할 패턴
HEADER_PATTERN = re.compile(r"^\d+장\s*·")  # "1장  ·  2026년의 교실"
FOOTER_PATTERN = re.compile(r"^(Living Textbook|교사 연수|예비교사 강의|\d+–\d+분)")
PAGE_NUM_PATTERN = re.compile(r"^\d{1,2}$")


def is_header_or_footer(text: str) -> bool:
    """헤더/푸터/페이지번호 판별."""
    t = text.strip()
    if not t:
        return True
    if HEADER_PATTERN.match(t):
        return True
    if FOOTER_PATTERN.match(t):
        return True
    if PAGE_NUM_PATTERN.match(t) and len(t) <= 3:
        return True
    return False


def classify_text(text: str) -> str:
    """텍스트가 헤더/푸터/번호인지 본문인지 분류."""
    return "skip" if is_header_or_footer(text) else "content"


def extract_slide_text(pptx_path: Path) -> list[dict]:
    """pptx 한 파일에서 슬라이드별 텍스트 추출."""
    prs = Presentation(pptx_path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        # 위치 정보로 정렬: 위→아래, 좌→우
        all_text = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            txt = shape.text_frame.text.strip()
            if not txt:
                continue
            try:
                top = shape.top if shape.top is not None else 0
                left = shape.left if shape.left is not None else 0
            except Exception:
                top, left = 0, 0
            all_text.append({"top": top, "left": left, "text": txt})
        all_text.sort(key=lambda d: (d["top"], d["left"]))

        # 헤더/푸터/페이지 번호 분리
        content_lines = []
        skip_lines = []
        for item in all_text:
            first_line = item["text"].split("\n")[0].strip()
            if classify_text(first_line) == "skip":
                skip_lines.append(item["text"])
            else:
                content_lines.append(item["text"])

        slides.append({
            "slide_no": i,
            "all_text_count": len(all_text),
            "content_text": "\n".join(content_lines),
            "skip_text": skip_lines,
        })
    return slides


def extract_manuscript(md_path: Path) -> dict:
    """마크다운 원고에서 챕터 메타 + 본문 추출."""
    raw = md_path.read_text(encoding="utf-8")
    lines = raw.split("\n")
    title = ""
    sections = []  # {"heading": ..., "paragraphs": [...]}
    current = None

    for line in lines:
        if line.startswith("# "):
            title = line[2:].strip()
            current = {"heading": title, "paragraphs": []}
            sections.append(current)
        elif line.startswith("## "):
            if current is None:
                current = {"heading": "", "paragraphs": []}
                sections.append(current)
            current = {"heading": line[3:].strip(), "paragraphs": []}
            sections.append(current)
        elif line.startswith("### "):
            if current is None:
                current = {"heading": "", "paragraphs": []}
                sections.append(current)
            current = {"heading": "## " + line[4:].strip(), "paragraphs": []}
            sections.append(current)
        elif line.strip() and current is not None:
            current["paragraphs"].append(line.strip())

    # 빈 섹션 제거
    sections = [s for s in sections if s["paragraphs"]]
    full_text = "\n\n".join(
        (s["heading"] + "\n" if s.get("heading") and not s["heading"].startswith("# ") else "")
        + "\n".join(s["paragraphs"])
        for s in sections
    )

    return {
        "title": title,
        "section_count": len(sections),
        "sections": sections,
        "full_text": full_text,
        "char_count": len(full_text),
    }


def build():
    OUT.mkdir(parents=True, exist_ok=True)
    manifest = json.loads(MANIFEST.read_text(encoding="utf-8"))

    slide_data = {}
    manuscript_data = {}
    chapter_data = []

    for ch in manifest["chapters"]:
        no = ch["chapter_no"]
        slide_path = Path(ch["slide_path"])
        md_path = Path(ch["manuscript_path"])

        slides = extract_slide_text(slide_path)
        slide_data[no] = {
            "file": str(slide_path),
            "chapter_no": no,
            "slide_count": len(slides),
            "slides": slides,
        }

        manuscript = extract_manuscript(md_path)
        manuscript_data[no] = {
            "file": str(md_path),
            "chapter_no": no,
            **manuscript,
        }

        chapter_data.append({
            "chapter_no": no,
            "section": ch["section"],
            "duration_min": ch["duration_min"],
            "slide_count": len(slides),
            "manuscript_chars": manuscript["char_count"],
            "manuscript_sections": manuscript["section_count"],
        })

    (OUT / "slide_texts.json").write_text(
        json.dumps(slide_data, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    (OUT / "manuscript_texts.json").write_text(
        json.dumps(manuscript_data, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    (OUT / "chapter_index.json").write_text(
        json.dumps(chapter_data, ensure_ascii=False, indent=2), encoding="utf-8"
    )

    print(f"✓ slide_texts.json     ({sum(len(d['slides']) for d in slide_data.values())} 슬라이드)")
    print(f"✓ manuscript_texts.json ({sum(d['section_count'] for d in manuscript_data.values())} 섹션)")
    print(f"✓ chapter_index.json    (17 챕터)")

    # 분량 진단
    print("\n--- 챕터별 분량 ---")
    for c in chapter_data:
        avg_chars_per_slide = c["manuscript_chars"] // max(c["slide_count"], 1)
        print(f"  [{c['chapter_no']}] 슬라이드 {c['slide_count']:>2}장 | 원고 {c['manuscript_chars']:>5}자 ({avg_chars_per_slide:>4}자/슬라이드) | {c['section']}")


if __name__ == "__main__":
    build()