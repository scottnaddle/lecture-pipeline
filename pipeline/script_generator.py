"""AI 스크립트 생성 (다중 프로바이더).

기존 generate_scripts.py를 ai_providers로 확장.
프로바이더는 환경변수 DEFAULT_AI_PROVIDER 또는 프로젝트 defaults.ai_provider로 선택.

사용법:
  python -m pipeline.script_generator --project-dir <path> --chapter 04
  python -m pipeline.script_generator --project-dir <path> --all
"""
from __future__ import annotations
import argparse
import json
import re
import sys
import time
from pathlib import Path
from typing import Optional

from .ai_providers import get_provider, detect_providers


SYSTEM_PROMPT = """당신은 한국어 강의 영상 나레이션을 작성하는 전문가입니다.
주어진 PPTX 슬라이드 텍스트와 원고를 바탕으로, 각 슬라이드당 60~80초 분량의
친근한 구어체 스크립트를 작성합니다.

규칙:
- 호칭은 "선생님" 사용
- 문장 호흡은 15~20어절
- 슬라이드 텍스트를 그대로 읽지 말고, 강의하듯이 풀어 설명
- 다음 슬라이드를 미리 언급하지 않음
- "여러분", "선생님" 같은 청취자 호명 사용
- 한 줄에 한 슬라이드: <번호>\t<나레이션>

출력 형식 (탭으로 구분):
1\t<슬라이드 1 나레이션>
2\t<슬라이드 2 나레이션>
...
"""


def read_slide_texts(pptx_path: Path) -> list[tuple[int, str]]:
    """PPTX에서 슬라이드 텍스트 추출 (python-pptx 사용)."""
    try:
        from pptx import Presentation
    except ImportError:
        print("python-pptx 미설치")
        return []
    if not pptx_path.exists():
        return []
    prs = Presentation(str(pptx_path))
    out = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs).strip()
                    if txt:
                        texts.append(txt)
        out.append((i, "\n".join(texts)))
    return out


def read_manuscript_excerpt(manuscript_path: Path, chapter_no: str, *, window: int = 2) -> str:
    """원고 파일에서 해당 챕터 주변 발췌. 챕터 경계를 ## Chapter N 같은 마커로 가정."""
    if not manuscript_path or not manuscript_path.exists():
        return ""
    text = manuscript_path.read_text(encoding="utf-8")
    # 챕터 마커 찾기 (## 1장, # Chapter 1, ## 01 등 다양한 형태)
    patterns = [
        rf"(?:#+\s*)?(?:Chapter|chapter|장|CHAPTER)\s*{int(chapter_no):d}\b",
        rf"(?:#+\s*)?(?:Chapter|chapter|장|CHAPTER)\s*{int(chapter_no):02d}\b",
        rf"(?:#+\s*){int(chapter_no):d}\s*장",
    ]
    boundaries = []
    for pat in patterns:
        for m in re.finditer(pat, text):
            boundaries.append(m.start())
    boundaries = sorted(set(boundaries))

    if not boundaries:
        # 마커 없으면 텍스트 일부 반환
        return text[:3000]

    # 가장 일치하는 위치
    start = boundaries[0]
    # 다음 마커 찾기
    nexts = [b for b in boundaries if b > start]
    end = nexts[0] if nexts else len(text)
    # 앞뒤 window 만큼 확장
    prevs = [b for b in boundaries if b < start]
    start = prevs[-1] if prevs else max(0, start - 200)
    end = min(len(text), end + 200)
    return text[start:end].strip()


def build_user_prompt(chapter_no: str, title: str, slide_texts: list[tuple[int, str]],
                     manuscript_excerpt: str) -> str:
    parts = [f"챕터: {chapter_no} - {title}\n"]
    parts.append("=== 슬라이드 텍스트 ===")
    for no, txt in slide_texts:
        parts.append(f"[슬라이드 {no}]\n{txt}\n")
    if manuscript_excerpt:
        parts.append("\n=== 관련 원고 발췌 ===\n" + manuscript_excerpt)
    parts.append("\n위 슬라이드 텍스트와 원고를 바탕으로, 각 슬라이드의 나레이션을 작성해주세요.")
    parts.append(f"출력은 {len(slide_texts)}줄, 각 줄은 '<번호>\\t<나레이션>' 형식입니다.")
    return "\n".join(parts)


def parse_script_output(text: str, expected_count: int) -> list[tuple[int, str]]:
    """AI 응답을 파싱. '번호\\t텍스트' 형식."""
    out = []
    for line in text.split("\n"):
        line = line.strip()
        if not line or "\t" not in line:
            continue
        # "1\t..." 또는 "1. ..." 등 다양한 형식 처리
        m = re.match(r"^(\d+)\s*[\.\)]\s*\t?\s*(.*)$", line) or re.match(r"^(\d+)\t(.*)$", line)
        if m:
            no = int(m.group(1))
            txt = m.group(2).strip()
            if 1 <= no <= expected_count and txt:
                out.append((no, txt))
    # 중복 제거 (같은 번호)
    seen = set()
    unique = []
    for no, txt in out:
        if no not in seen:
            seen.add(no)
            unique.append((no, txt))
    return sorted(unique, key=lambda x: x[0])


def generate_chapter_script(
    project_dir: Path, chapter_no: str,
    *, provider_name: Optional[str] = None,
    temperature: float = 0.7,
) -> int:
    """한 챕터의 스크립트 생성. 성공 시 슬라이드 수 반환."""
    proj_path = project_dir / "project.json"
    state = json.loads(proj_path.read_text(encoding="utf-8"))

    # 챕터 정보
    if chapter_no not in state["chapters"]:
        print(f"✗ ch{chapter_no} 없음")
        return 0
    ch = state["chapters"][chapter_no]
    title = ch.get("title", "")
    if not title:
        # manifest에서 slide_path로 추출
        mfst = json.loads((project_dir / "00_매니페스트" / "manifest.json").read_text(encoding="utf-8"))
        for mc in mfst.get("chapters", []):
            if mc["chapter_no"] == chapter_no:
                title = Path(mc["slide_path"]).stem.split("_", 1)[-1] if mc.get("slide_path") else ""
                break

    # PPTX 슬라이드 텍스트
    slide_path = None
    mfst = json.loads((project_dir / "00_매니페스트" / "manifest.json").read_text(encoding="utf-8"))
    for mc in mfst.get("chapters", []):
        if mc["chapter_no"] == chapter_no:
            slide_path = Path(mc["slide_path"]) if mc.get("slide_path") else None
            break
    slide_texts = read_slide_texts(slide_path) if slide_path else []
    if not slide_texts:
        print(f"✗ ch{chapter_no}: PPTX 슬라이드 텍스트 없음")
        return 0

    # 원고 발췌
    manuscript_path = Path(state.get("manuscript_path", ""))
    manuscript_excerpt = read_manuscript_excerpt(manuscript_path, chapter_no)

    # AI 호출
    provider = get_provider(provider_name or state.get("defaults", {}).get("ai_provider"))
    print(f"[{provider.name}/{provider.model}] ch{chapter_no} ({len(slide_texts)}장) 생성 중...")
    t0 = time.time()
    try:
        user_prompt = build_user_prompt(chapter_no, title, slide_texts, manuscript_excerpt)
        response = provider.chat(SYSTEM_PROMPT, user_prompt, temperature=temperature, max_tokens=8192)
    except Exception as e:
        print(f"✗ AI 호출 실패: {e}")
        return 0
    dt = time.time() - t0

    # 파싱
    parsed = parse_script_output(response, len(slide_texts))
    if len(parsed) < len(slide_texts) * 0.8:  # 80% 미만이면 부분만
        print(f"  ⚠ {len(parsed)}/{len(slide_texts)}장만 파싱됨 (전체 {dt:.0f}초)")

    # 누락된 슬라이드 채우기 (있으면)
    parsed_dict = dict(parsed)
    for no, _ in slide_texts:
        if no not in parsed_dict:
            parsed_dict[no] = f"(슬라이드 {no} 자동 생성 실패 - 직접 작성 필요)"

    # 파일 저장
    out_path = project_dir / "01_스크립트" / "scripts" / f"ch{chapter_no}.txt"
    out_path.parent.mkdir(parents=True, exist_ok=True)
    lines = [
        f"# 챕터 {chapter_no} 자동 생성 스크립트 (AI: {provider.name})",
        f"# 형식: <slide_no>\\t<text>",
        f"# 생성: {time.strftime('%Y-%m-%d %H:%M:%S')}",
        ""
    ]
    for no, txt in sorted(parsed_dict.items()):
        lines.append(f"{no}\t{txt}")
    out_path.write_text("\n".join(lines) + "\n", encoding="utf-8")
    print(f"  ✓ {len(parsed_dict)}/{len(slide_texts)}장 저장: {out_path} ({dt:.0f}초)")
    return len(parsed_dict)


def main():
    p = argparse.ArgumentParser(description="AI 스크립트 생성")
    p.add_argument("--project-dir", required=True)
    p.add_argument("--chapter", help="특정 챕터")
    p.add_argument("--all", action="store_true", help="모든 챕터")
    p.add_argument("--provider", help="AI 프로바이더 (openai|anthropic|zhipu|moonshot|MiniMax)")
    p.add_argument("--temperature", type=float, default=0.7)
    args = p.parse_args()

    pd = Path(args.project_dir).resolve()
    if not (pd / "project.json").exists():
        print(f"project.json 없음: {pd}")
        sys.exit(1)

    state = json.loads((pd / "project.json").read_text(encoding="utf-8"))

    if args.chapter:
        generate_chapter_script(pd, args.chapter, provider_name=args.provider, temperature=args.temperature)
    elif args.all:
        for ch_no in sorted(state["chapters"].keys()):
            generate_chapter_script(pd, ch_no, provider_name=args.provider, temperature=args.temperature)
    else:
        print("--chapter NN 또는 --all 필요")
        sys.exit(1)


if __name__ == "__main__":
    main()
